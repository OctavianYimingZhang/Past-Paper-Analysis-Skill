#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import math
import re
import subprocess
import sys
import zipfile
from collections import Counter, defaultdict
from dataclasses import asdict, dataclass, field
from pathlib import Path
from statistics import mean
from typing import Iterable

import fitz
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation


SKILL_ROOT = Path(__file__).resolve().parent.parent
SCRIPT_ROOT = SKILL_ROOT / "scripts"
REFERENCE_ROOT = SKILL_ROOT / "references"
DEFAULT_OUTPUT_ROOT = Path.cwd() / "output" / "past-paper-knowledge-point-analysis"
DEFAULT_CACHE_ROOT = Path.cwd() / "tmp" / "past-paper-knowledge-point-analysis"
OCR_SOURCE = SCRIPT_ROOT / "vision_ocr.swift"
OCR_BINARY = DEFAULT_CACHE_ROOT / "vision_ocr"

HEADER_FILL = PatternFill("solid", fgColor="1F4E78")
HEADER_FONT = Font(color="FFFFFF", bold=True)
WRAP = Alignment(vertical="top", wrap_text=True)

STOPWORDS = {
    "a",
    "about",
    "after",
    "all",
    "an",
    "and",
    "are",
    "as",
    "at",
    "be",
    "because",
    "before",
    "between",
    "both",
    "by",
    "can",
    "correct",
    "do",
    "does",
    "each",
    "either",
    "english",
    "explanation",
    "for",
    "from",
    "has",
    "have",
    "how",
    "if",
    "in",
    "into",
    "is",
    "it",
    "its",
    "key",
    "many",
    "more",
    "most",
    "not",
    "of",
    "on",
    "only",
    "option",
    "or",
    "question",
    "so",
    "some",
    "than",
    "that",
    "the",
    "their",
    "therefore",
    "these",
    "they",
    "this",
    "those",
    "to",
    "use",
    "using",
    "what",
    "when",
    "which",
    "with",
    "wrong",
}

BOILERPLATE_SNIPPETS = (
    "the university of manchester",
    "multiple choice question",
    "answer sheet instructions",
    "electronic calculators may be used",
    "do not remove this paper",
    "downloaded by",
    "lomoarcpsd",
    "page ",
    "section a",
)

ANSWER_VARIANTS = (
    ("answer_english_option_text", re.compile(r"Answer\s*\(English option text\)\s*:\s*(.+?)(?=\n(?:Explanation|Correct answer|Answer)|\Z)", re.I | re.S)),
    ("answer_english_only", re.compile(r"Answer\s*\(English only\)\s*:\s*(.+?)(?=\n(?:Explanation|Correct answer|Answer)|\Z)", re.I | re.S)),
    ("answer_english_option_content_only", re.compile(r"Answer\s*\(English option content only\)\s*:\s*(.+?)(?=\n(?:Explanation|Correct answer|Answer)|\Z)", re.I | re.S)),
    ("correct_answer", re.compile(r"Correct answer\s*:\s*(.+?)(?=\n(?:Explanation|Correct answer|Answer)|\Z)", re.I | re.S)),
)

EXPLANATION_VARIANTS = (
    ("explanation_english", re.compile(r"Explanation\s*\(English\)\s*:\s*(.+)", re.I | re.S)),
    ("explanation_en", re.compile(r"Explanation\s*\(EN\)\s*:\s*(.+)", re.I | re.S)),
    ("explanation_plain", re.compile(r"Explanation\s*:\s*(.+)", re.I | re.S)),
)

RANGE_CONTEXT_RE = re.compile(
    r"Questions?\s+(\d+)\s*(?:to|-|–|and)\s*(\d+)\b(.*?)(?=(?:\[\d+\]|(?:\n\s*\d+\.\s+)|\Z))",
    re.I | re.S,
)


@dataclass
class PaperSpec:
    year: str
    pdf: Path
    expected_questions: int
    role: str = "formal"
    question_numbers: list[int] | None = None
    parser_strategy: str | None = None


@dataclass
class AnswerKeySpec:
    year: str
    docx: Path


@dataclass
class NotesConfig:
    ocr_fallback_mode: str
    min_detected_lectures: int
    lecture_header_patterns: list[str]


@dataclass
class CourseSpec:
    course_id: str
    course_name: str
    slides_dir: Path | None
    notes_pdf: Path | None
    formal_years: list[str]
    paper_skip_pages: int
    notes: NotesConfig
    papers: list[PaperSpec]
    answer_keys: list[AnswerKeySpec]
    preset_id: str | None
    parser_hints: dict[str, object]
    output_language: str
    slides_manifest: list[dict[str, object]]
    manual_overrides: dict[str, dict[str, object]]


@dataclass(frozen=True)
class Preset:
    preset_id: str
    preferred_source: str
    benchmark_kind: str


PRESETS = {
    "biochemistry": Preset("biochemistry", "slides", "biochemistry"),
    "from-molecules-to-cells": Preset("from-molecules-to-cells", "slides", "from-molecules-to-cells"),
    "drugs": Preset("drugs", "slides", "drugs"),
    "excitable-cells": Preset("excitable-cells", "slides", "excitable-cells"),
}


@dataclass
class OCRItemSummary:
    index: int
    text: str
    avg_confidence: float


@dataclass
class OCRDocumentSummary:
    path: str
    items: list[OCRItemSummary]


@dataclass
class Lecture:
    lecture_id: str
    lecture_number: int
    title: str
    start_page: int
    end_page: int
    raw_text: str
    raw_candidates: list[str]
    optimized_points: list["OptimizedPoint"] = field(default_factory=list)

    @property
    def page_range(self) -> str:
        return f"{self.start_page}-{self.end_page}"


@dataclass
class OptimizedPoint:
    point_id: str
    lecture_id: str
    label: str
    source_candidates: list[str]
    tokens: set[str]


@dataclass
class ReviewIssue:
    severity: str
    kind: str
    source: str
    year: str | None
    question_number: int | None
    detail: str


@dataclass
class QuestionRecord:
    year: str
    question_number: int
    question_id: str
    stem: str
    full_text: str
    source_pages: str
    answer_text: str
    answer_variant: str | None
    explanation_text: str
    explanation_variant: str | None
    primary_lecture_id: str | None
    primary_point_id: str | None
    secondary_point_ids: list[str]
    lecture_score: float
    point_score: float
    confidence: float
    review_flags: list[str]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--spec", required=True, help="Path to a course spec JSON file")
    parser.add_argument("--output-root", default=str(DEFAULT_OUTPUT_ROOT))
    parser.add_argument("--cache-root", default=str(DEFAULT_CACHE_ROOT))
    parser.add_argument("--force-note-ocr", action="store_true")
    return parser.parse_args()


def default_notes_config() -> NotesConfig:
    return NotesConfig(
        ocr_fallback_mode="text_first",
        min_detected_lectures=6,
        lecture_header_patterns=[
            r"Lecture\s*(\d{1,2})",
            r"\bL(\d{1,2})\b",
        ],
    )


def load_course_spec(path: Path) -> CourseSpec:
    raw = json.loads(path.read_text(encoding="utf-8"))
    if raw.get("output_language", "en") != "en":
        raise ValueError("Only English output is currently supported.")
    if not raw.get("slides_dir") and not raw.get("notes_pdf"):
        raise ValueError("Course spec requires at least one of slides_dir or notes_pdf.")
    notes = NotesConfig(**raw["notes"]) if raw.get("notes") else default_notes_config()
    papers = [
        PaperSpec(
            year=str(item["year"]),
            pdf=Path(item["pdf"]),
            expected_questions=int(item["expected_questions"]),
            role=str(item.get("role", "formal")),
            question_numbers=[int(value) for value in item["question_numbers"]] if item.get("question_numbers") else None,
            parser_strategy=item.get("parser_strategy"),
        )
        for item in raw["papers"]
    ]
    answer_keys = [
        AnswerKeySpec(year=str(item["year"]), docx=Path(item["docx"]))
        for item in raw.get("answer_keys", [])
    ]
    formal_years = list(raw.get("formal_years") or [paper.year for paper in papers if paper.role == "formal"])
    return CourseSpec(
        course_id=raw["course_id"],
        course_name=raw["course_name"],
        slides_dir=Path(raw["slides_dir"]) if raw.get("slides_dir") else None,
        notes_pdf=Path(raw["notes_pdf"]) if raw.get("notes_pdf") else None,
        formal_years=formal_years,
        paper_skip_pages=int(raw.get("paper_skip_pages", 0)),
        notes=notes,
        papers=papers,
        answer_keys=answer_keys,
        preset_id=raw.get("preset_id"),
        parser_hints=dict(raw.get("parser_hints", {})),
        output_language=str(raw.get("output_language", "en")),
        slides_manifest=list(raw.get("slides_manifest", [])),
        manual_overrides=dict(raw.get("manual_overrides", {})),
    )


def resolve_preset(spec: CourseSpec) -> Preset | None:
    if spec.preset_id and spec.preset_id in PRESETS:
        return PRESETS[spec.preset_id]
    course_slug = make_slug(spec.course_name)
    return PRESETS.get(course_slug)


def ensure_dirs(*paths: Path) -> None:
    for path in paths:
        path.mkdir(parents=True, exist_ok=True)


def normalize_whitespace(text: str) -> str:
    text = text.replace("\xa0", " ").replace("\x0b", " ")
    text = text.replace("\u2013", "-").replace("\u2014", "-")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def flatten_text(text: str) -> str:
    return re.sub(r"\s+", " ", normalize_whitespace(text)).strip()


def strip_cjk(text: str) -> str:
    return "".join(char for char in text if not ("\u4e00" <= char <= "\u9fff"))


def english_output_text(text: str) -> str:
    return normalize_whitespace(strip_cjk(text)).strip()


def normalize_for_search(text: str) -> str:
    text = flatten_text(text).lower()
    replacements = {
        "α": " alpha ",
        "β": " beta ",
        "γ": " gamma ",
        "δ": " delta ",
        "μ": " micro ",
        "+": " plus ",
        "/": " ",
    }
    for source, target in replacements.items():
        text = text.replace(source, target)
    text = re.sub(r"[^0-9a-z\- ]+", " ", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def tokenize(text: str) -> list[str]:
    tokens = []
    for token in normalize_for_search(text).split():
        if token in STOPWORDS:
            continue
        if len(token) <= 1 and not token.isdigit():
            continue
        tokens.append(token)
    return tokens


def make_slug(text: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", normalize_for_search(text))
    return slug.strip("-")


def compile_ocr_binary(cache_root: Path) -> Path:
    ensure_dirs(cache_root)
    binary = cache_root / "vision_ocr"
    if binary.exists() and binary.stat().st_mtime >= OCR_SOURCE.stat().st_mtime:
        return binary
    subprocess.run(["swiftc", str(OCR_SOURCE), "-o", str(binary)], check=True)
    return binary


def run_ocr_pdf(pdf_path: Path, binary: Path, cache_root: Path) -> OCRDocumentSummary:
    cache_path = cache_root / f"{make_slug(pdf_path.stem)}.ocr.pages.json"
    if cache_path.exists():
        return load_ocr_summary(cache_path)
    result = subprocess.run([str(binary), "--pdf", str(pdf_path)], capture_output=True, text=True, check=True)
    cache_path.write_text(result.stdout, encoding="utf-8")
    return load_ocr_summary(cache_path)


def render_pdf_to_images(pdf_path: Path, output_dir: Path, scale: float = 2.5) -> list[Path]:
    ensure_dirs(output_dir)
    document = fitz.open(pdf_path)
    matrix = fitz.Matrix(scale, scale)
    paths = []
    for index, page in enumerate(document, start=1):
        target = output_dir / f"{make_slug(pdf_path.stem)}-p{index:02d}.png"
        if not target.exists():
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            pix.save(target)
        paths.append(target)
    return paths


def run_ocr_rendered_pdf(pdf_path: Path, binary: Path, cache_root: Path, scale: float = 2.5) -> OCRDocumentSummary:
    image_dir = cache_root / f"{make_slug(pdf_path.stem)}_rendered_pages"
    image_paths = render_pdf_to_images(pdf_path, image_dir, scale=scale)
    return run_ocr_images(image_paths, f"{make_slug(pdf_path.stem)}.rendered", binary, cache_root)


def extract_docx_images(docx_path: Path, output_dir: Path) -> list[Path]:
    ensure_dirs(output_dir)
    with zipfile.ZipFile(docx_path) as archive:
        names = [
            name
            for name in archive.namelist()
            if name.startswith("word/media/") and name.lower().endswith((".png", ".jpg", ".jpeg"))
        ]
        names.sort(key=lambda name: int(re.search(r"(\d+)", Path(name).stem).group(1)))
        paths = []
        for name in names:
            target = output_dir / Path(name).name
            if not target.exists():
                target.write_bytes(archive.read(name))
            paths.append(target)
    return paths


def run_ocr_images(image_paths: list[Path], cache_name: str, binary: Path, cache_root: Path) -> OCRDocumentSummary:
    cache_path = cache_root / f"{cache_name}.ocr.images.json"
    if cache_path.exists():
        return load_ocr_summary(cache_path)
    cmd = [str(binary), "--images", *[str(path) for path in image_paths]]
    result = subprocess.run(cmd, capture_output=True, text=True, check=True)
    cache_path.write_text(result.stdout, encoding="utf-8")
    return load_ocr_summary(cache_path)


def load_ocr_summary(path: Path) -> OCRDocumentSummary:
    raw = json.loads(path.read_text(encoding="utf-8"))
    items = []
    for item in raw["items"]:
        confidences = [float(line["confidence"]) for line in item.get("lines", [])]
        items.append(
            OCRItemSummary(
                index=int(item["index"]),
                text=normalize_whitespace(item.get("text", "")),
                avg_confidence=round(mean(confidences), 4) if confidences else 0.0,
            )
        )
    return OCRDocumentSummary(path=raw.get("path", ""), items=items)


def extract_page_payloads(pdf_path: Path, start_page: int = 0) -> list[dict[str, object]]:
    document = fitz.open(pdf_path)
    payloads: list[dict[str, object]] = []
    for page_index in range(start_page, document.page_count):
        page = document.load_page(page_index)
        text = normalize_whitespace(page.get_text("text"))
        lines = [normalize_whitespace(line) for line in page.get_text("text").splitlines() if normalize_whitespace(line)]
        payloads.append({"index": page_index + 1, "text": text, "lines": lines})
    return payloads


def infer_lecture_number_from_filename(file_name: str, fallback: int) -> int:
    patterns = [
        r"lecture\s*0*(\d{1,2})\b",
        r"\bl0*(\d{1,2})\b",
        r"^0*(\d{1,2})[\s._-]",
    ]
    lowered = file_name.lower()
    for pattern in patterns:
        match = re.search(pattern, lowered)
        if match:
            return int(match.group(1))
    return fallback


def infer_lecture_title_from_filename(file_name: str) -> str:
    stem = Path(file_name).stem
    stem = re.sub(r"^\d+(?:\.\d+)?[\s._-]*", "", stem)
    stem = re.sub(r"\blecture\b", "", stem, flags=re.I)
    stem = re.sub(r"\bslides?\b", "", stem, flags=re.I)
    stem = re.sub(r"\bno movies\b", "", stem, flags=re.I)
    stem = re.sub(r"\bcopy\b", "", stem, flags=re.I)
    stem = re.sub(r"[_-]+", " ", stem)
    stem = re.sub(r"\s+", " ", stem).strip(" .-_")
    return english_output_text(stem[:140] if stem else Path(file_name).stem)


def extract_pptx_text(path: Path) -> tuple[str, list[str]]:
    presentation = Presentation(str(path))
    slide_lines: list[str] = []
    titles: list[str] = []
    for slide in presentation.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = normalize_whitespace(shape.text or "")
                if text:
                    texts.append(text)
        if texts:
            titles.append(texts[0])
            slide_lines.extend(texts)
    return "\n".join(slide_lines), titles


def section_slides_into_lectures(spec: CourseSpec, review_issues: list[ReviewIssue]) -> tuple[list[Lecture], str]:
    if spec.slides_dir is None:
        return [], "none"
    if not spec.slides_dir.exists():
        review_issues.append(
            ReviewIssue(
                severity="warning",
                kind="slides_dir_missing",
                source=display_source_label(spec.slides_dir),
                year=None,
                question_number=None,
                detail="slides_dir was provided in the spec but does not exist.",
            )
        )
        return [], "missing"
    lectures: list[Lecture] = []
    manifest_by_name = {str(item["file_name"]): item for item in spec.slides_manifest}
    slide_files = sorted(spec.slides_dir.glob("*.pptx"))
    for index, slide_path in enumerate(slide_files, start=1):
        raw_text, titles = extract_pptx_text(slide_path)
        if slide_path.name in manifest_by_name:
            manifest = manifest_by_name[slide_path.name]
            lecture_id = str(manifest["lecture_id"])
            lecture_number = int(manifest.get("lecture_number", re.search(r"(\d+)", lecture_id).group(1)))
            title = english_output_text(str(manifest["lecture_title"]))
        else:
            lecture_number = infer_lecture_number_from_filename(slide_path.name, index)
            lecture_id = f"L{lecture_number:02d}"
            title = english_output_text(titles[0]) if titles else infer_lecture_title_from_filename(slide_path.name)
        raw_candidates = extract_candidate_knowledge_points(raw_text.splitlines(), title)
        lecture = Lecture(
            lecture_id=lecture_id,
            lecture_number=lecture_number,
            title=title,
            start_page=1,
            end_page=max(1, len(titles)),
            raw_text=raw_text,
            raw_candidates=raw_candidates,
        )
        lecture.optimized_points = optimize_candidates(lecture.lecture_id, lecture.raw_candidates)
        lectures.append(lecture)
    lectures.sort(key=lambda item: (item.lecture_number, item.lecture_id))
    return lectures, "slides"


def detect_lecture_number(text: str, patterns: list[str]) -> int | None:
    for pattern in patterns:
        match = re.search(pattern, text, re.I)
        if match:
            try:
                return int(match.group(1))
            except ValueError:
                continue
    return None


def extract_lecture_title(text: str, number: int) -> str:
    for line in normalize_whitespace(text).splitlines():
        flattened = english_output_text(flatten_text(line))
        if not flattened:
            continue
        if re.search(rf"lecture\s*{number}\b", flattened, re.I):
            return flattened[:140]
    return f"Lecture {number}"


def section_notes_into_lectures(spec: CourseSpec, binary: Path, cache_root: Path, force_note_ocr: bool, review_issues: list[ReviewIssue]) -> tuple[list[Lecture], str]:
    extraction_mode = "text"
    if spec.notes_pdf is None:
        return [], "none"
    payloads = extract_page_payloads(spec.notes_pdf)
    lecture_hits = sum(1 for payload in payloads if detect_lecture_number(str(payload["text"]), spec.notes.lecture_header_patterns) is not None)
    use_ocr = force_note_ocr or (spec.notes.ocr_fallback_mode == "ocr_only") or lecture_hits < spec.notes.min_detected_lectures
    if use_ocr:
        ocr_doc = run_ocr_pdf(spec.notes_pdf, binary, cache_root)
        payloads = [{"index": item.index, "text": item.text, "lines": item.text.splitlines()} for item in ocr_doc.items]
        extraction_mode = "ocr"
        lecture_hits = sum(1 for payload in payloads if detect_lecture_number(str(payload["text"]), spec.notes.lecture_header_patterns) is not None)
    if lecture_hits < spec.notes.min_detected_lectures:
        review_issues.append(
            ReviewIssue(
                severity="warning",
                kind="notes_lecture_detection",
                source=display_source_label(spec.notes_pdf),
                year=None,
                question_number=None,
                detail=f"Detected only {lecture_hits} lecture headers in {extraction_mode} mode.",
            )
        )
    lectures: list[Lecture] = []
    current_number: int | None = None
    current_start = 0
    current_texts: list[str] = []
    current_lines: list[str] = []
    current_title = ""
    for payload in payloads:
        text = str(payload["text"])
        number = detect_lecture_number(text, spec.notes.lecture_header_patterns)
        if number is not None and number != current_number:
            if current_number is not None and current_texts:
                lecture = build_lecture(current_number, current_title, current_start, int(payload["index"]) - 1, current_texts, current_lines)
                lectures.append(lecture)
            current_number = number
            current_start = int(payload["index"])
            current_texts = [text]
            current_lines = list(payload["lines"])
            current_title = extract_lecture_title(text, number)
        elif current_number is not None:
            current_texts.append(text)
            current_lines.extend(payload["lines"])
    if current_number is not None and current_texts:
        lecture = build_lecture(current_number, current_title, current_start, int(payloads[-1]["index"]), current_texts, current_lines)
        lectures.append(lecture)
    lectures.sort(key=lambda item: item.lecture_number)
    return lectures, extraction_mode


def extract_lectures(
    spec: CourseSpec,
    binary: Path,
    cache_root: Path,
    force_note_ocr: bool,
    review_issues: list[ReviewIssue],
) -> tuple[list[Lecture], str]:
    preset = resolve_preset(spec)
    prefer_slides = bool(spec.slides_dir) and (preset is not None and preset.preferred_source == "slides")
    if prefer_slides:
        lectures, mode = section_slides_into_lectures(spec, review_issues)
        if lectures:
            return lectures, mode
    if spec.slides_dir:
        lectures, mode = section_slides_into_lectures(spec, review_issues)
        if lectures and spec.notes_pdf is None:
            return lectures, mode
    lectures, mode = section_notes_into_lectures(spec, binary, cache_root, force_note_ocr, review_issues)
    if lectures:
        return lectures, mode
    if spec.slides_dir:
        lectures, slide_mode = section_slides_into_lectures(spec, review_issues)
        if lectures:
            return lectures, slide_mode
    return [], "none"


def build_lecture(number: int, title: str, start_page: int, end_page: int, page_texts: list[str], lines: list[str]) -> Lecture:
    raw_text = "\n".join(page_texts)
    raw_candidates = extract_candidate_knowledge_points(lines, title)
    lecture = Lecture(
        lecture_id=f"L{number:02d}",
        lecture_number=number,
        title=title,
        start_page=start_page,
        end_page=end_page,
        raw_text=raw_text,
        raw_candidates=raw_candidates,
    )
    lecture.optimized_points = optimize_candidates(lecture.lecture_id, lecture.raw_candidates)
    return lecture


def extract_candidate_knowledge_points(lines: list[str], title: str) -> list[str]:
    counter: Counter[str] = Counter()
    best_form: dict[str, str] = {}
    for line in lines:
        candidate = english_output_text(flatten_text(line))
        if not is_candidate_line(candidate, title):
            continue
        normalized = normalize_for_search(candidate)
        counter[normalized] += 1
        stored = best_form.get(normalized)
        if stored is None or len(candidate) > len(stored):
            best_form[normalized] = candidate
    if not counter:
        return []
    ranked = sorted(counter.items(), key=lambda item: (-item[1], -candidate_quality(best_form[item[0]]), best_form[item[0]].lower()))
    return [best_form[key] for key, _count in ranked[:40]]


def is_candidate_line(line: str, title: str) -> bool:
    lower = line.lower()
    if len(line) < 18 or len(line) > 160:
        return False
    if line == title:
        return False
    if any(snippet in lower for snippet in BOILERPLATE_SNIPPETS):
        return False
    if lower.startswith(("lecture ", "module ", "biol", "page ")):
        return False
    if re.search(r"\bquestion\b|\banswer\b", lower):
        return False
    alpha = sum(1 for char in line if char.isalpha())
    if alpha < max(8, int(len(line) * 0.35)):
        return False
    token_count = len(tokenize(line))
    if token_count < 3 or token_count > 20:
        return False
    return True


def candidate_quality(text: str) -> int:
    return len(tokenize(text)) * 10 + min(len(text), 120)


def optimize_candidates(lecture_id: str, raw_candidates: list[str]) -> list[OptimizedPoint]:
    clusters: list[dict[str, object]] = []
    for candidate in raw_candidates:
        tokens = set(tokenize(candidate))
        if not tokens:
            continue
        assigned = False
        for cluster in clusters:
            cluster_tokens = cluster["tokens"]
            similarity = token_jaccard(tokens, cluster_tokens)
            if similarity >= 0.72 or normalize_for_search(candidate) in normalize_for_search(str(cluster["label"])) or normalize_for_search(str(cluster["label"])) in normalize_for_search(candidate):
                cluster["candidates"].append(candidate)
                cluster["tokens"] = set(cluster_tokens) | tokens
                if candidate_quality(candidate) > candidate_quality(str(cluster["label"])):
                    cluster["label"] = candidate
                assigned = True
                break
        if not assigned:
            clusters.append({"label": candidate, "tokens": tokens, "candidates": [candidate]})
    points = []
    for index, cluster in enumerate(clusters[:20], start=1):
        points.append(
            OptimizedPoint(
                point_id=f"{lecture_id}-T{index:02d}",
                lecture_id=lecture_id,
                label=str(cluster["label"]),
                source_candidates=list(cluster["candidates"]),
                tokens=set(cluster["tokens"]),
            )
        )
    return points


def token_jaccard(left: set[str], right: set[str]) -> float:
    if not left or not right:
        return 0.0
    return len(left & right) / len(left | right)


def build_combined_text(payloads: list[dict[str, object]]) -> str:
    return "\n\n".join(f"[[PAGE {payload['index']}]]\n{payload['text']}" for payload in payloads if payload["text"])


def locate_markers(text: str, question_numbers: list[int], marker_type: str) -> list[tuple[int, int, int, str]]:
    markers = []
    search_from = 0
    for number in question_numbers:
        candidates = []
        for label, pattern in marker_patterns(number, marker_type):
            match = re.search(pattern, text[search_from:], re.I | re.M)
            if match:
                candidates.append((search_from + match.start(), search_from + match.end(), label))
        if not candidates:
            return []
        start, end, label = min(candidates, key=lambda item: item[0])
        markers.append((number, start, end, label))
        search_from = end
    return markers


def marker_patterns(number: int, marker_type: str) -> list[tuple[str, str]]:
    if marker_type == "paper":
        return [
            ("dot", rf"(?<![A-Za-z0-9]){number}\.\s+"),
            ("bracket", rf"\[{number}\]"),
            ("question_word", rf"Question\s+{number}\b"),
        ]
    return [
        ("qline", rf"^\s*Q\s*{number}\s*(?=$|\n)"),
        ("nline", rf"^\s*{number}\s*(?=$|\n)"),
        ("qinline", rf"\bQ\s*{number}\b"),
    ]


def strip_marker(block: str, number: int, label: str) -> str:
    if label == "dot":
        return re.sub(rf"^\s*{number}\.\s*", "", block, count=1)
    if label == "bracket":
        return re.sub(rf"^\s*\[{number}\]\s*", "", block, count=1)
    if label == "question_word":
        return re.sub(rf"^\s*Question\s+{number}\s*", "", block, count=1, flags=re.I)
    if label == "qline":
        return re.sub(rf"^\s*Q\s*{number}\s*", "", block, count=1, flags=re.I)
    if label == "nline":
        return re.sub(rf"^\s*{number}\s*", "", block, count=1)
    return re.sub(rf"\bQ\s*{number}\b", "", block, count=1, flags=re.I)


def extract_range_contexts(text: str) -> list[tuple[int, int, str]]:
    contexts = []
    for start, end, context in RANGE_CONTEXT_RE.findall(text):
        contexts.append((int(start), int(end), flatten_text(context)))
    return contexts


def build_question_blocks(text: str, question_numbers: list[int], marker_type: str) -> list[dict[str, object]]:
    markers = locate_markers(text, question_numbers, marker_type)
    if not markers:
        return []
    contexts = extract_range_contexts(text) if marker_type == "paper" else []
    blocks = []
    for index, (number, start, end, label) in enumerate(markers):
        block_end = markers[index + 1][1] if index + 1 < len(markers) else len(text)
        block = text[start:block_end]
        block = strip_marker(block, number, label)
        pages = ",".join(sorted(set(re.findall(r"\[\[PAGE (\d+)\]\]", block))))
        block = re.sub(r"\[\[PAGE \d+\]\]\n?", "", block)
        context_text = ""
        for start_number, end_number, context in contexts:
            if start_number <= number <= end_number:
                context_text = context
                break
        combined = normalize_whitespace(" ".join(part for part in [context_text, block] if part))
        blocks.append({"number": number, "text": combined, "pages": pages, "label": label, "context": context_text})
    return blocks


def stem_from_block(block_text: str) -> str:
    option_match = re.search(r"(?:^|\n)\s*A[\.\)]\s+", block_text)
    if not option_match:
        return flatten_text(block_text)
    return flatten_text(block_text[: option_match.start()])


def extract_questions_from_paper(
    paper: PaperSpec,
    paper_skip_pages: int,
    binary: Path,
    cache_root: Path,
    review_issues: list[ReviewIssue],
) -> tuple[list[dict[str, object]], str]:
    question_numbers = paper.question_numbers or list(range(1, paper.expected_questions + 1))
    payloads = extract_page_payloads(paper.pdf, start_page=paper_skip_pages)
    method = "text"
    blocks = build_question_blocks(build_combined_text(payloads), question_numbers, "paper")
    if len(blocks) != len(question_numbers):
        payloads = extract_page_payloads(paper.pdf, start_page=0)
        blocks = build_question_blocks(build_combined_text(payloads), question_numbers, "paper")
        method = "text_full"
    if len(blocks) != len(question_numbers):
        ocr_doc = run_ocr_pdf(paper.pdf, binary, cache_root)
        ocr_payloads = [{"index": item.index, "text": item.text} for item in ocr_doc.items[paper_skip_pages:]]
        blocks = build_question_blocks(build_combined_text(ocr_payloads), question_numbers, "paper")
        method = "ocr"
    if len(blocks) != len(question_numbers):
        ocr_payloads = [{"index": item.index, "text": item.text} for item in ocr_doc.items]
        blocks = build_question_blocks(build_combined_text(ocr_payloads), question_numbers, "paper")
        method = "ocr_full"
    if len(blocks) != len(question_numbers):
        rendered_ocr = run_ocr_rendered_pdf(paper.pdf, binary, cache_root)
        rendered_payloads = [{"index": item.index, "text": item.text} for item in rendered_ocr.items[paper_skip_pages:]]
        blocks = build_question_blocks(build_combined_text(rendered_payloads), question_numbers, "paper")
        method = "ocr_rendered"
    if len(blocks) != len(question_numbers):
        rendered_payloads = [{"index": item.index, "text": item.text} for item in rendered_ocr.items]
        blocks = build_question_blocks(build_combined_text(rendered_payloads), question_numbers, "paper")
        method = "ocr_rendered_full"
    if len(blocks) != len(question_numbers):
        review_issues.append(
            ReviewIssue(
                severity="error",
                kind="paper_question_count",
                source=display_source_label(paper.pdf),
                year=paper.year,
                question_number=None,
                detail=f"Recovered {len(blocks)} of {len(question_numbers)} questions using {method}.",
            )
        )
    questions = []
    for block in blocks:
        questions.append(
            {
                "year": paper.year,
                "question_number": int(block["number"]),
                "question_id": f"{paper.year}-Q{int(block['number']):02d}",
                "stem": english_output_text(stem_from_block(str(block["text"]))),
                "full_text": english_output_text(str(block["text"])),
                "source_pages": str(block["pages"]),
            }
        )
    return questions, method


def parse_answer_block(block: str) -> tuple[str, str, str | None, str | None]:
    answer_text = ""
    answer_variant = None
    explanation_text = ""
    explanation_variant = None
    for label, pattern in ANSWER_VARIANTS:
        match = pattern.search(block)
        if match:
            answer_text = clean_answer_field(match.group(1))
            answer_variant = label
            break
    for label, pattern in EXPLANATION_VARIANTS:
        match = pattern.search(block)
        if match:
            explanation_text = clean_explanation_field(match.group(1))
            explanation_variant = label
            break
    return answer_text, explanation_text, answer_variant, explanation_variant


def clean_answer_field(text: str) -> str:
    text = english_output_text(text)
    text = re.sub(r"^[\-\u2022]+", "", text).strip()
    return text[:320]


def clean_explanation_field(text: str) -> str:
    text = re.sub(r"(?:^|\n)\s*#.*", "", text)
    text = english_output_text(text)
    return text[:4000]


def extract_answers_from_key(
    answer_key: AnswerKeySpec,
    question_numbers: list[int],
    binary: Path,
    cache_root: Path,
    review_issues: list[ReviewIssue],
) -> tuple[dict[int, dict[str, object]], dict[str, int]]:
    image_dir = cache_root / f"{make_slug(answer_key.docx.stem)}_images"
    image_paths = extract_docx_images(answer_key.docx, image_dir)
    ocr_doc = run_ocr_images(image_paths, make_slug(answer_key.docx.stem), binary, cache_root)
    combined_text = "\n\n".join(item.text for item in ocr_doc.items if item.text)
    blocks = build_question_blocks(combined_text, question_numbers, "answer")
    variant_counts: dict[str, int] = defaultdict(int)
    answers: dict[int, dict[str, object]] = {}
    if len(blocks) != len(question_numbers):
        review_issues.append(
            ReviewIssue(
                severity="error",
                kind="answer_key_question_count",
                source=display_source_label(answer_key.docx),
                year=answer_key.year,
                question_number=None,
                detail=f"Recovered {len(blocks)} of {len(question_numbers)} answer blocks from OCR.",
            )
        )
    course_avg_confidence = mean([item.avg_confidence for item in ocr_doc.items if item.text]) if ocr_doc.items else 0.0
    if course_avg_confidence < 0.68:
        review_issues.append(
            ReviewIssue(
                severity="warning",
                kind="answer_key_low_ocr_confidence",
                source=display_source_label(answer_key.docx),
                year=answer_key.year,
                question_number=None,
                detail=f"Average OCR confidence was {course_avg_confidence:.3f}.",
            )
        )
    for block in blocks:
        answer_text, explanation_text, answer_variant, explanation_variant = parse_answer_block(str(block["text"]))
        if answer_variant:
            variant_counts[answer_variant] += 1
        if explanation_variant:
            variant_counts[explanation_variant] += 1
        if not answer_text or not explanation_text:
            review_issues.append(
                ReviewIssue(
                    severity="warning",
                    kind="answer_key_parse",
                    source=display_source_label(answer_key.docx),
                    year=answer_key.year,
                    question_number=int(block["number"]),
                    detail="Missing answer text or explanation text after OCR parsing.",
                )
            )
        answers[int(block["number"])] = {
            "answer_text": answer_text,
            "answer_variant": answer_variant,
            "explanation_text": explanation_text,
            "explanation_variant": explanation_variant,
        }
    return answers, dict(variant_counts)


def build_idf(documents: Iterable[Iterable[str]]) -> dict[str, float]:
    doc_list = [set(doc) for doc in documents if doc]
    total = len(doc_list)
    if total == 0:
        return {}
    df: Counter[str] = Counter()
    for doc in doc_list:
        df.update(doc)
    return {token: math.log((1 + total) / (1 + count)) + 1.0 for token, count in df.items()}


def score_tokens(query_tokens: list[str], target_tokens: set[str], idf: dict[str, float]) -> float:
    if not query_tokens or not target_tokens:
        return 0.0
    shared = set(query_tokens) & target_tokens
    return round(sum(idf.get(token, 1.0) for token in shared), 4)


def confidence_from_scores(scores: list[float]) -> float:
    if not scores or scores[0] <= 0:
        return 0.0
    top = scores[0]
    second = scores[1] if len(scores) > 1 else 0.0
    margin = (top - second) / top if top else 0.0
    scale = min(top / 8.0, 1.0)
    return round((margin * 0.65) + (scale * 0.35), 3)


def apply_manual_override(question: QuestionRecord, override: dict[str, object]) -> QuestionRecord:
    if "lecture_id" in override:
        question.primary_lecture_id = str(override["lecture_id"])
    if "point_id" in override:
        question.primary_point_id = str(override["point_id"])
    if "secondary_point_ids" in override:
        question.secondary_point_ids = [str(value) for value in override["secondary_point_ids"]]
    if "notes" in override:
        question.review_flags.append(str(override["notes"]))
    question.confidence = max(question.confidence, 0.95)
    return question


def map_questions_to_points(
    spec: CourseSpec,
    lectures: list[Lecture],
    paper_questions: list[dict[str, object]],
    answer_map: dict[str, dict[int, dict[str, object]]],
    review_issues: list[ReviewIssue],
) -> list[QuestionRecord]:
    lecture_tokens = {lecture.lecture_id: set(tokenize(lecture.raw_text + " " + " ".join(lecture.raw_candidates))) for lecture in lectures}
    point_lookup = {point.point_id: point for lecture in lectures for point in lecture.optimized_points}
    lecture_idf = build_idf(lecture_tokens.values())
    point_idf = build_idf(point.tokens for point in point_lookup.values())
    questions: list[QuestionRecord] = []
    for item in paper_questions:
        year = str(item["year"])
        answer_payload = answer_map.get(year, {}).get(int(item["question_number"]), {})
        query_text = normalize_whitespace(" ".join(part for part in [str(item["stem"]), str(answer_payload.get("answer_text", "")), str(answer_payload.get("explanation_text", ""))] if part))
        query_tokens = tokenize(query_text)
        lecture_scored = []
        for lecture in lectures:
            score = score_tokens(query_tokens, lecture_tokens.get(lecture.lecture_id, set()), lecture_idf)
            lecture_scored.append((lecture.lecture_id, score))
        lecture_scored.sort(key=lambda entry: entry[1], reverse=True)
        primary_lecture_id = lecture_scored[0][0] if lecture_scored and lecture_scored[0][1] > 0 else None
        lecture_confidence = confidence_from_scores([score for _lecture_id, score in lecture_scored[:3]])
        lecture_score = lecture_scored[0][1] if lecture_scored else 0.0
        point_scored = []
        if primary_lecture_id:
            lecture = next(lecture for lecture in lectures if lecture.lecture_id == primary_lecture_id)
            for point in lecture.optimized_points:
                score = score_tokens(query_tokens, point.tokens, point_idf)
                point_scored.append((point.point_id, score))
            point_scored.sort(key=lambda entry: entry[1], reverse=True)
        primary_point_id = point_scored[0][0] if point_scored and point_scored[0][1] > 0 else None
        point_score = point_scored[0][1] if point_scored else 0.0
        point_confidence = confidence_from_scores([score for _point_id, score in point_scored[:3]])
        confidence = round((lecture_confidence * 0.55) + (point_confidence * 0.45), 3)
        secondary_point_ids = [
            point_id
            for point_id, score in point_scored[1:3]
            if point_score > 0 and score >= point_score * 0.85 and score > 0
        ]
        review_flags = []
        if lecture_score <= 0:
            review_flags.append("weak_lecture_mapping")
        if point_score <= 0:
            review_flags.append("weak_point_mapping")
        if confidence < 0.45:
            review_flags.append("low_mapping_confidence")
        question = QuestionRecord(
            year=year,
            question_number=int(item["question_number"]),
            question_id=str(item["question_id"]),
            stem=str(item["stem"]),
            full_text=str(item["full_text"]),
            source_pages=str(item["source_pages"]),
            answer_text=str(answer_payload.get("answer_text", "")),
            answer_variant=answer_payload.get("answer_variant"),
            explanation_text=str(answer_payload.get("explanation_text", "")),
            explanation_variant=answer_payload.get("explanation_variant"),
            primary_lecture_id=primary_lecture_id,
            primary_point_id=primary_point_id,
            secondary_point_ids=secondary_point_ids,
            lecture_score=lecture_score,
            point_score=point_score,
            confidence=confidence,
            review_flags=review_flags,
        )
        override_key = f"{year}-Q{question.question_number:02d}"
        if override_key in spec.manual_overrides:
            question = apply_manual_override(question, spec.manual_overrides[override_key])
        if review_flags:
            review_issues.append(
                ReviewIssue(
                    severity="warning",
                    kind="question_mapping",
                    source=question.question_id,
                    year=question.year,
                    question_number=question.question_number,
                    detail=", ".join(review_flags),
                )
            )
        questions.append(question)
    return questions


def recover_missing_questions_from_answer_keys(
    paper_questions: list[dict[str, object]],
    answer_map: dict[str, dict[int, dict[str, object]]],
    question_numbers_by_year: dict[str, list[int]],
    review_issues: list[ReviewIssue],
) -> list[dict[str, object]]:
    recovered = list(paper_questions)
    existing = {(str(item["year"]), int(item["question_number"])) for item in paper_questions}
    for year, question_numbers in question_numbers_by_year.items():
        for question_number in question_numbers:
            key = (year, question_number)
            if key in existing:
                continue
            if question_number not in answer_map.get(year, {}):
                continue
            recovered.append(
                {
                    "year": year,
                    "question_number": question_number,
                    "question_id": f"{year}-Q{question_number:02d}",
                    "stem": f"Question {question_number}",
                    "full_text": f"Question {question_number}",
                    "source_pages": "",
                }
            )
            review_issues.append(
                ReviewIssue(
                    severity="warning",
                    kind="paper_answer_key_recovery",
                    source=f"{year}-Q{question_number:02d}",
                    year=year,
                    question_number=question_number,
                    detail="Paper question block was missing; recovered a placeholder from the answer key so the topic can still be mapped.",
                )
            )
    recovered.sort(key=lambda item: (str(item["year"]), int(item["question_number"])))
    return recovered


def retention_band(years_present: int, total_years: int) -> str:
    if years_present <= 0:
        return "Not tested"
    if total_years <= 0:
        return "Not tested"
    ratio = years_present / total_years
    if ratio >= 0.75:
        return "Anchor"
    if ratio >= 0.50:
        return "Core"
    if years_present > 1:
        return "Recurring"
    return "One-off"


def build_topic_rows(lectures: list[Lecture], questions: list[QuestionRecord], formal_years: list[str]) -> list[dict[str, object]]:
    year_counts: dict[str, Counter[str]] = defaultdict(Counter)
    total_counts: Counter[str] = Counter()
    for question in questions:
        if not question.primary_point_id:
            continue
        total_counts[question.primary_point_id] += 1
        year_counts[question.primary_point_id][question.year] += 1
    all_points = [point for lecture in lectures for point in lecture.optimized_points]
    sorted_points = sorted(all_points, key=lambda point: (-total_counts[point.point_id], point.point_id))
    hit_points = [point for point in sorted_points if total_counts[point.point_id] > 0]
    hotness_rank = {point.point_id: index for index, point in enumerate(hit_points, start=1)}
    total_questions = sum(total_counts.values())
    rows = []
    for point in sorted_points:
        counts = year_counts[point.point_id]
        years_present = sum(1 for year in formal_years if counts.get(year, 0) > 0)
        retention_percent = round((years_present / len(formal_years)) * 100, 1) if formal_years else 0.0
        row = {
            "topic_id": point.point_id,
            "lecture_id": point.lecture_id,
            "optimized_topic": point.label,
            "raw_question_hits": total_counts[point.point_id],
            "question_share_percent": round((total_counts[point.point_id] / total_questions) * 100, 2) if total_questions else 0.0,
            "hotness_rank": hotness_rank.get(point.point_id, ""),
            "years_present": years_present,
            "retention_fraction": f"{years_present}/{len(formal_years)}",
            "retention_percent": retention_percent,
            "meets_50": retention_percent >= 50.0 if total_counts[point.point_id] else False,
            "meets_75": retention_percent >= 75.0 if total_counts[point.point_id] else False,
            "retention_band": retention_band(years_present, len(formal_years)),
        }
        for year in formal_years:
            row[year] = counts.get(year, 0)
        rows.append(row)
    return rows


def build_method_rows(spec: CourseSpec, lecture_mode: str, paper_methods: dict[str, str], answer_variant_counts: dict[str, int], review_issues: list[ReviewIssue]) -> list[list[object]]:
    paper_roles = ", ".join(f"{paper.year}:{paper.role}" for paper in spec.papers)
    rows = [
        ["Course", spec.course_name],
        ["Course ID", spec.course_id],
        ["Preset ID", spec.preset_id or "generic"],
        ["Output language", spec.output_language],
        ["Slides source", display_source_label(spec.slides_dir) if spec.slides_dir else "None"],
        ["Notes source", display_source_label(spec.notes_pdf) if spec.notes_pdf else "None"],
        ["Lecture extraction mode", lecture_mode],
        ["Formal years", ", ".join(spec.formal_years)],
        ["Paper roles", paper_roles],
        ["Hotness definition", "Primary mapped topic raw_question_hits, question_share_percent, hotness_rank, and per-year counts."],
        ["Retention definition", "years_present / total_years with explicit meets_50, meets_75, and retention_band fields."],
        ["Retention bands", "Anchor >= 75%; Core >= 50% and < 75%; Recurring > 1 year but < 50%; One-off = 1 year; Not tested = 0 hits."],
        ["Review gate", "Weak OCR, weak question splitting, missing question numbers, weak lecture/topic mapping, or undercount vs expected_questions go to Review_Queue."],
        ["Paper extraction methods", ", ".join(f"{year}:{method}" for year, method in sorted(paper_methods.items()))],
        ["Answer parser variants seen", ", ".join(f"{key}:{value}" for key, value in sorted(answer_variant_counts.items())) or "No answer-key OCR used"],
        ["Review issue count", len(review_issues)],
    ]
    return rows


def write_sheet(ws, headers: list[str], rows: list[list[object]]) -> None:
    ws.append([sanitize_excel_value(value) for value in headers])
    for cell in ws[1]:
        cell.fill = HEADER_FILL
        cell.font = HEADER_FONT
        cell.alignment = WRAP
    ws.freeze_panes = "A2"
    ws.auto_filter.ref = f"A1:{get_column_letter(len(headers))}1"
    for row in rows:
        ws.append([sanitize_excel_value(value) for value in row])
    for row in ws.iter_rows():
        for cell in row:
            cell.alignment = WRAP
    autofit(ws)


def autofit(ws) -> None:
    for column_cells in ws.columns:
        width = 12
        for cell in column_cells:
            if cell.value is None:
                continue
            width = max(width, min(len(str(cell.value)) + 2, 60))
        ws.column_dimensions[get_column_letter(column_cells[0].column)].width = width


def sanitize_excel_value(value: object) -> object:
    if not isinstance(value, str):
        return value
    filtered = []
    for char in value:
        code = ord(char)
        if char in {"\t", "\n", "\r"}:
            filtered.append(char)
            continue
        if 0x20 <= code <= 0xD7FF or 0xE000 <= code <= 0xFFFD or 0x10000 <= code <= 0x10FFFF:
            filtered.append(char)
        else:
            filtered.append(" ")
    return "".join(filtered)


def display_source_label(value: str | Path | None) -> str:
    if value is None:
        return ""
    if isinstance(value, Path):
        return value.name or str(value)
    if "/" in value:
        return Path(value).name or value
    return value


def render_workbook(
    spec: CourseSpec,
    lectures: list[Lecture],
    questions: list[QuestionRecord],
    topic_rows: list[dict[str, object]],
    review_issues: list[ReviewIssue],
    output_root: Path,
    lecture_mode: str,
    paper_methods: dict[str, str],
    answer_variant_counts: dict[str, int],
) -> Path:
    ensure_dirs(output_root)
    output_path = output_root / f"{spec.course_id}.xlsx"
    wb = Workbook()
    default = wb.active
    wb.remove(default)

    ws = wb.create_sheet("Method")
    write_sheet(ws, ["Field", "Value"], build_method_rows(spec, lecture_mode, paper_methods, answer_variant_counts, review_issues))

    lecture_rows = []
    for lecture in lectures:
        lecture_rows.append(
            [
                lecture.lecture_id,
                lecture.lecture_number,
                lecture.title,
                lecture.page_range,
                "\n".join(f"- {item}" for item in lecture.raw_candidates[:20]),
                "\n".join(f"- {point.point_id}: {point.label}" for point in lecture.optimized_points),
            ]
        )
    ws = wb.create_sheet("Lecture_Knowledge_Map")
    write_sheet(ws, ["Lecture_ID", "Lecture_Number", "Lecture_Title", "Page_Range", "Initial_Candidates", "Optimized_Points"], lecture_rows)

    question_rows = []
    for question in questions:
        question_rows.append(
            [
                question.year,
                question.question_number,
                question.question_id,
                question.source_pages,
                question.stem,
                question.answer_text,
                question.explanation_text,
                question.primary_lecture_id or "",
                question.primary_point_id or "",
                ", ".join(question.secondary_point_ids),
                question.lecture_score,
                question.point_score,
                question.confidence,
                ", ".join(question.review_flags),
            ]
        )
    ws = wb.create_sheet("Question_Mapping")
    write_sheet(
        ws,
        [
            "Year",
            "Question_Number",
            "Question_ID",
            "Source_Pages",
            "Stem",
            "Answer_Text",
            "Explanation_Text",
            "Primary_Lecture_ID",
            "Primary_Point_ID",
            "Secondary_Point_IDs",
            "Lecture_Score",
            "Point_Score",
            "Confidence",
            "Review_Flags",
        ],
        question_rows,
    )

    topic_headers = ["Topic_ID", "Lecture_ID", "Optimized_Topic", "Raw_Question_Hits", "Question_Share_Percent", "Hotness_Rank", *spec.formal_years, "Years_Present", "Retention_Fraction", "Retention_Percent", "Meets_50", "Meets_75", "Retention_Band"]
    topic_sheet_rows = []
    for row in topic_rows:
        topic_sheet_rows.append(
            [
                row["topic_id"],
                row["lecture_id"],
                row["optimized_topic"],
                row["raw_question_hits"],
                row["question_share_percent"],
                row["hotness_rank"],
                *[row[year] for year in spec.formal_years],
                row["years_present"],
                row["retention_fraction"],
                row["retention_percent"],
                row["meets_50"],
                row["meets_75"],
                row["retention_band"],
            ]
        )
    ws = wb.create_sheet("Topic_Frequency")
    write_sheet(ws, topic_headers, topic_sheet_rows)

    ws = wb.create_sheet("Year_Topic_Matrix")
    write_sheet(ws, topic_headers, topic_sheet_rows)

    retention_rows = [row for row in topic_sheet_rows if row[3] > 0]
    retention_rows.sort(key=lambda item: (band_sort_key(str(item[-1])), -(item[3] or 0), str(item[2]).lower()))
    ws = wb.create_sheet("Retention_Bands")
    write_sheet(ws, topic_headers, retention_rows)

    review_rows = [
        [issue.severity, issue.kind, issue.source, issue.year or "", issue.question_number or "", issue.detail]
        for issue in review_issues
    ]
    ws = wb.create_sheet("Review_Queue")
    write_sheet(ws, ["Severity", "Kind", "Source", "Year", "Question_Number", "Detail"], review_rows)

    wb.save(output_path)
    return output_path


def band_sort_key(label: str) -> int:
    order = {"Anchor": 0, "Core": 1, "Recurring": 2, "One-off": 3, "Not tested": 4}
    return order.get(label, 9)


def render_json(
    spec: CourseSpec,
    lectures: list[Lecture],
    questions: list[QuestionRecord],
    topic_rows: list[dict[str, object]],
    review_issues: list[ReviewIssue],
    output_root: Path,
    workbook_path: Path,
    markdown_path: Path,
    lecture_mode: str,
    paper_methods: dict[str, str],
    paper_question_counts: dict[str, int],
    answer_key_question_counts: dict[str, int],
    answer_variant_counts: dict[str, int],
) -> Path:
    ensure_dirs(output_root)
    output_path = output_root / f"{spec.course_id}.analysis.json"
    summary = {
        "course_id": spec.course_id,
        "course_name": spec.course_name,
        "preset_id": spec.preset_id or "generic",
        "output_language": spec.output_language,
        "lecture_extraction_mode": lecture_mode,
        "expected_questions_total": sum(paper.expected_questions for paper in spec.papers),
        "recovered_questions_total": len(questions),
        "paper_roles": {paper.year: paper.role for paper in spec.papers},
        "paper_question_counts": paper_question_counts,
        "answer_key_question_counts": answer_key_question_counts,
        "paper_extraction_methods": paper_methods,
        "answer_parser_variant_counts": answer_variant_counts,
        "review_issue_count": len(review_issues),
        "output_files": {
            "workbook": str(workbook_path),
            "markdown": str(markdown_path),
        },
    }
    payload = {
        "summary": summary,
        "lectures": [
            {
                "lecture_id": lecture.lecture_id,
                "lecture_number": lecture.lecture_number,
                "title": lecture.title,
                "page_range": lecture.page_range,
                "raw_candidates": lecture.raw_candidates,
                "optimized_points": [asdict(point) | {"tokens": sorted(point.tokens)} for point in lecture.optimized_points],
            }
            for lecture in lectures
        ],
        "questions": [asdict(question) for question in questions],
        "topics": topic_rows,
        "review_queue": [asdict(issue) for issue in review_issues],
    }
    output_path.write_text(json.dumps(payload, indent=2, ensure_ascii=False), encoding="utf-8")
    summary_path = output_root / f"{spec.course_id}.summary.json"
    summary_path.write_text(json.dumps(summary, indent=2, ensure_ascii=False), encoding="utf-8")
    return output_path


def render_markdown(
    spec: CourseSpec,
    topic_rows: list[dict[str, object]],
    review_issues: list[ReviewIssue],
    output_root: Path,
    paper_question_counts: dict[str, int],
    answer_key_question_counts: dict[str, int],
) -> Path:
    ensure_dirs(output_root)
    output_path = output_root / f"{spec.course_id}.summary.md"
    anchors = [row for row in topic_rows if row["retention_band"] == "Anchor" and row["raw_question_hits"] > 0][:10]
    cores = [row for row in topic_rows if row["retention_band"] == "Core" and row["raw_question_hits"] > 0][:10]
    review_preview = review_issues[:20]
    lines = [
        f"# {spec.course_name}",
        "",
        f"- Formal years: {', '.join(spec.formal_years)}",
        f"- Papers analyzed: {', '.join(f'{paper.year} ({paper.role})' for paper in spec.papers)}",
        f"- Expected questions: {sum(paper.expected_questions for paper in spec.papers)}",
        f"- Recovered questions: {sum(paper_question_counts.values())}",
        f"- Review issues: {len(review_issues)}",
        "",
        "## Paper Counts",
        "",
    ]
    for paper in spec.papers:
        year = paper.year
        lines.append(f"- {year}: paper={paper_question_counts.get(year, 0)}, answer_key={answer_key_question_counts.get(year, 0)}")
    lines.extend(["", "## Anchor Topics", ""])
    if anchors:
        for row in anchors:
            lines.append(f"- {row['topic_id']} | {row['optimized_topic']} | hits={row['raw_question_hits']} | retention={row['retention_fraction']} ({row['retention_percent']}%)")
    else:
        lines.append("- None")
    lines.extend(["", "## Core Topics", ""])
    if cores:
        for row in cores:
            lines.append(f"- {row['topic_id']} | {row['optimized_topic']} | hits={row['raw_question_hits']} | retention={row['retention_fraction']} ({row['retention_percent']}%)")
    else:
        lines.append("- None")
    lines.extend(["", "## Review Queue Preview", ""])
    if review_preview:
        for issue in review_preview:
            label = f"{issue.year}-Q{issue.question_number:02d}" if issue.year and issue.question_number else issue.source
            lines.append(f"- {issue.kind} | {label} | {issue.detail}")
    else:
        lines.append("- None")
    output_path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return output_path


def run_analysis(spec: CourseSpec, output_root: Path, cache_root: Path, force_note_ocr: bool) -> dict[str, object]:
    ensure_dirs(output_root, cache_root)
    review_issues: list[ReviewIssue] = []
    binary = compile_ocr_binary(cache_root)
    lectures, lecture_mode = extract_lectures(spec, binary, cache_root, force_note_ocr, review_issues)
    if not lectures:
        raise RuntimeError(f"No lectures detected for course {spec.course_id}")
    paper_questions: list[dict[str, object]] = []
    paper_methods: dict[str, str] = {}
    paper_question_counts: dict[str, int] = {}
    for paper in spec.papers:
        extracted, method = extract_questions_from_paper(paper, spec.paper_skip_pages, binary, cache_root, review_issues)
        paper_questions.extend(extracted)
        paper_methods[paper.year] = method
        paper_question_counts[paper.year] = len(extracted)
    answer_map: dict[str, dict[int, dict[str, object]]] = {}
    answer_key_question_counts: dict[str, int] = {}
    answer_variant_counts: Counter[str] = Counter()
    question_numbers_by_year = {paper.year: (paper.question_numbers or list(range(1, paper.expected_questions + 1))) for paper in spec.papers}
    if spec.answer_keys:
        for answer_key in spec.answer_keys:
            answers, variant_counts = extract_answers_from_key(answer_key, question_numbers_by_year[answer_key.year], binary, cache_root, review_issues)
            answer_map[answer_key.year] = answers
            answer_key_question_counts[answer_key.year] = len(answers)
            answer_variant_counts.update(variant_counts)
    else:
        review_issues.append(
            ReviewIssue(
                severity="warning",
                kind="answer_key_missing",
                source=spec.course_id,
                year=None,
                question_number=None,
                detail="No answer keys were supplied. Mapping quality will rely on slides/notes and question stems only.",
            )
        )
    paper_questions = recover_missing_questions_from_answer_keys(paper_questions, answer_map, question_numbers_by_year, review_issues)
    questions = map_questions_to_points(spec, lectures, paper_questions, answer_map, review_issues)
    retention_years = spec.formal_years or [paper.year for paper in spec.papers if paper.role == "formal"]
    topic_rows = build_topic_rows(lectures, [question for question in questions if question.year in retention_years or any(p.year == question.year for p in spec.papers)], retention_years)
    workbook_path = render_workbook(spec, lectures, questions, topic_rows, review_issues, output_root, lecture_mode, paper_methods, dict(answer_variant_counts))
    markdown_path = render_markdown(spec, topic_rows, review_issues, output_root, paper_question_counts, answer_key_question_counts)
    json_path = render_json(spec, lectures, questions, topic_rows, review_issues, output_root, workbook_path, markdown_path, lecture_mode, paper_methods, paper_question_counts, answer_key_question_counts, dict(answer_variant_counts))
    return {
        "workbook": workbook_path,
        "markdown": markdown_path,
        "json": json_path,
        "review_count": len(review_issues),
    }


def main() -> None:
    args = parse_args()
    spec_path = Path(args.spec).expanduser().resolve()
    output_root = Path(args.output_root).expanduser().resolve()
    cache_root = Path(args.cache_root).expanduser().resolve()
    spec = load_course_spec(spec_path)
    result = run_analysis(spec, output_root / spec.course_id, cache_root / spec.course_id, args.force_note_ocr)
    print(json.dumps({key: str(value) if isinstance(value, Path) else value for key, value in result.items()}, indent=2))


if __name__ == "__main__":
    try:
        main()
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        raise
